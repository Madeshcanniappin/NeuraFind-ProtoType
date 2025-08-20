import os
import datetime
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def read_file(filepath):
    text = ""
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == ".txt":
            with open(filepath, "r", errors="ignore") as f:
                text = f.read()
        elif ext == ".pdf":
            reader = PdfReader(filepath)
            for page in reader.pages:
                text += page.extract_text() or ""
        elif ext == ".docx":
            doc = Document(filepath)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == ".pptx":
            pres = Presentation(filepath)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            text = f"[Unsupported format: {filepath}]"
    except Exception as e:
        text = f"[Error reading {filepath}: {str(e)}]"

    return text


def extract_metadata(filepath):
    stats = os.stat(filepath)
    return {
        "size": stats.st_size,
        "created": datetime.datetime.fromtimestamp(stats.st_ctime),
        "modified": datetime.datetime.fromtimestamp(stats.st_mtime),
        "type": os.path.splitext(filepath)[1]
    }


def log_results(query, semantic_results, metadata_results):
    with open("results/search_log.txt", "a") as f:
        f.write(f"\nQuery: {query}\n")
        f.write("Semantic Results:\n")
        for r in semantic_results:
            f.write(f"  - {r['file']} (score: {r['score']:.4f})\n")
        f.write("Metadata Results:\n")
        for r in metadata_results:
            f.write(f"  - {r['file']} (match: {r['match']})\n")
